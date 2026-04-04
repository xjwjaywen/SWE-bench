"""生成100封假邮件数据和对应附件文件 — 金融信托公司场景。"""

import csv
import json
import os
import random
import zipfile
from datetime import datetime, timedelta
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation as PptxPresentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas as pdf_canvas

# 路径配置
BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
EMAILS_DIR = DATA_DIR / "emails"
ATTACHMENTS_DIR = DATA_DIR / "attachments"

# ============================================================
# 人员和公司数据 — 华信金融信托有限公司
# ============================================================
PEOPLE = [
    {"name": "张明远", "email": "zhangmy@huaxintrust.com", "dept": "信托业务部"},
    {"name": "李文博", "email": "liwb@huaxintrust.com", "dept": "财富管理部"},
    {"name": "王建华", "email": "wangjh@huaxintrust.com", "dept": "风控合规部"},
    {"name": "赵雪梅", "email": "zhaoxm@huaxintrust.com", "dept": "法务部"},
    {"name": "孙浩然", "email": "sunhr@huaxintrust.com", "dept": "资产管理部"},
    {"name": "周丽萍", "email": "zhoulp@huaxintrust.com", "dept": "运营管理部"},
    {"name": "吴志强", "email": "wuzq@huaxintrust.com", "dept": "投资研究部"},
    {"name": "郑晓东", "email": "zhengxd@huaxintrust.com", "dept": "合规审计部"},
    {"name": "钱海燕", "email": "qianhy@huaxintrust.com", "dept": "客户服务部"},
    {"name": "陈思远", "email": "chensy@huaxintrust.com", "dept": "信托业务部"},
    {"name": "林国栋", "email": "lingd@huaxintrust.com", "dept": "投资研究部"},
    {"name": "黄永华", "email": "huangyh@huaxintrust.com", "dept": "财务部"},
    {"name": "刘德伟", "email": "liudw@hengfundinvest.com", "dept": "恒丰投资"},
    {"name": "杨芳", "email": "yangf@longtaiwealth.com", "dept": "龙泰财富"},
    {"name": "马骏", "email": "maj@zhongyuancapital.com", "dept": "中原资本"},
]


def pick_sender():
    return random.choice(PEOPLE)


def pick_recipients(exclude_email, count=None):
    count = count or random.randint(1, 3)
    pool = [p for p in PEOPLE if p["email"] != exclude_email]
    return random.sample(pool, min(count, len(pool)))


def random_date(start_year=2024):
    start = datetime(start_year, 1, 1)
    delta = timedelta(days=random.randint(0, 450))
    return (start + delta).strftime("%Y-%m-%d %H:%M:%S")


# ============================================================
# 附件生成器（不变）
# ============================================================

def make_pdf(filepath: str, title: str, content: str):
    c = pdf_canvas.Canvas(filepath, pagesize=A4)
    c.setFont("Helvetica", 16)
    c.drawString(72, 750, title)
    c.setFont("Helvetica", 11)
    y = 720
    for line in content.split("\n"):
        if y < 72:
            c.showPage()
            c.setFont("Helvetica", 11)
            y = 750
        c.drawString(72, y, line)
        y -= 16
    c.save()


def make_docx(filepath: str, title: str, content: str):
    doc = Document()
    doc.add_heading(title, level=1)
    for para in content.split("\n\n"):
        doc.add_paragraph(para)
    doc.save(filepath)


def make_xlsx(filepath: str, title: str, rows: list[list]):
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]
    for row in rows:
        ws.append(row)
    wb.save(filepath)


def make_pptx(filepath: str, title: str, slides_content: list[tuple[str, str]]):
    prs = PptxPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Huaxin Trust Co., Ltd."
    for slide_title, slide_body in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = slide_body
    prs.save(filepath)


def make_csv_file(filepath: str, headers: list[str], rows: list[list]):
    with open(filepath, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerow(headers)
        writer.writerows(rows)


def make_txt(filepath: str, content: str):
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)


def make_png(filepath: str, text: str, size=(400, 300)):
    img = Image.new("RGB", size, color=(240, 240, 245))
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, size[0] - 10, size[1] - 10], outline=(100, 100, 200), width=2)
    draw.text((20, 20), text, fill=(50, 50, 50))
    img.save(filepath)


def make_jpg(filepath: str, text: str, size=(400, 300)):
    img = Image.new("RGB", size, color=(255, 250, 240))
    draw = ImageDraw.Draw(img)
    draw.text((20, 20), text, fill=(80, 40, 40))
    img.save(filepath, "JPEG")


def make_zip(filepath: str, files_content: dict[str, str]):
    with zipfile.ZipFile(filepath, "w") as zf:
        for name, content in files_content.items():
            zf.writestr(name, content)


# ============================================================
# 邮件模板 — 8个金融业务场景
# ============================================================

def gen_trust_emails(start_id: int) -> list[dict]:
    """信托业务 — 15封：信托计划设立、管理、清算"""
    emails = []
    trust_names = [
        "华信·稳健增长1号集合资金信托计划",
        "华信·产业投资2号单一资金信托计划",
        "华信·城市发展3号集合资金信托计划",
        "华信·医疗健康专项信托计划",
        "华信·科技创新股权投资信托计划",
    ]
    topics = [
        ("信托计划设立方案", "setup"),
        ("信托计划推介材料", "pitch"),
        ("信托合同定稿", "contract"),
        ("受托人报告", "trustee_report"),
        ("信托财产管理报告", "asset_report"),
        ("信托计划季度运营报告", "quarterly"),
        ("信托资金投放通知", "funding"),
        ("信托收益分配方案", "distribution"),
        ("信托计划延期公告", "extension"),
        ("信托清算报告", "liquidation"),
        ("信托项目尽调报告", "due_diligence"),
        ("信托计划变更公告", "amendment"),
        ("委托人大会通知", "meeting"),
        ("信托登记信息", "registration"),
        ("信托项目风险评估", "risk_assessment"),
    ]

    for i in range(15):
        eid = f"email_{start_id + i:03d}"
        trust = trust_names[i % len(trust_names)]
        topic_name, topic_key = topics[i]
        sender = pick_sender()
        recipients = pick_recipients(sender["email"])
        date = random_date()
        subject = f"[信托业务] {trust} - {topic_name}"

        # 预先生成共享数据（邮件正文和附件使用相同数值，保持一致）
        trust_scale_wan = random.randint(5000, 50000)
        trust_duration_str = random.choice(['12个月', '18个月', '24个月', '36个月'])
        trust_rate_pct = round(random.uniform(5.0, 8.5), 2)
        trust_invest_dir = random.choice(['基础设施建设', '房地产开发', '工商企业贷款', '证券投资', '股权投资'])
        trust_risk_measure = random.choice(['土地抵押+应收账款质押', '股权质押+担保', '优先/劣后结构化设计'])
        due_diligence_rating = random.choice(['AA+', 'AA', 'AA-', 'A+'])
        due_diligence_growth = random.randint(8, 25)
        due_diligence_debt = random.randint(40, 70)
        due_diligence_collateral = random.randint(2, 20)
        due_diligence_ltv = random.randint(40, 65)
        dist_amount_wan = random.randint(500, 5000)
        dist_rate_pct = round(random.uniform(5.5, 8.0), 2)
        liquidation_earnings = random.randint(800, 8000)
        liquidation_cost = random.randint(10, 50)
        liquidation_total = random.randint(10000, 50000)
        att_loan_assets = random.randint(3000, 30000)
        att_cash_assets = random.randint(500, 3000)
        att_interest_recv = random.randint(100, 800)

        # 根据不同主题生成不同邮件正文
        if "设立" in topic_name:
            body_text = (
                f"各位同事好，\n\n"
                f"关于{trust}的设立方案，现将相关材料发送如下：\n\n"
                f"一、信托计划基本要素：\n"
                f"- 信托规模：{trust_scale_wan}万元\n"
                f"- 信托期限：{trust_duration_str}\n"
                f"- 预期收益率：{trust_rate_pct}%/年\n"
                f"- 投资方向：{trust_invest_dir}\n"
                f"- 风控措施：{trust_risk_measure}\n\n"
                f"二、信托计划设立流程：\n"
                f"1. 项目立项审批\n"
                f"2. 尽职调查\n"
                f"3. 信托文件起草与审批\n"
                f"4. 监管报备\n"
                f"5. 产品推介与资金募集\n"
                f"6. 信托计划成立\n\n"
                f"请审阅附件中的详细方案，如有意见请于本周五前反馈。\n\n{sender['name']}\n信托业务部"
            )
        elif "尽调" in topic_name:
            body_text = (
                f"各位好，\n\n"
                f"附件为{trust}的尽职调查报告，主要调查结果如下：\n\n"
                f"1. 融资方主体信用评级：{due_diligence_rating}\n"
                f"2. 融资方近三年营收复合增长率：{due_diligence_growth}%\n"
                f"3. 资产负债率：{due_diligence_debt}%\n"
                f"4. 抵押物评估价值：{due_diligence_collateral}亿元\n"
                f"5. 抵押率：{due_diligence_ltv}%\n\n"
                f"风险提示：\n"
                f"- {random.choice(['融资方所在行业存在周期性波动', '抵押物变现存在一定流动性风险', '区域经济下行压力需关注'])}\n"
                f"- 建议增加{random.choice(['保证金账户监管', '现金流归集', '信息披露频率'])}要求\n\n"
                f"{sender['name']}"
            )
        elif "收益分配" in topic_name:
            body_text = (
                f"各位委托人/同事好，\n\n"
                f"关于{trust}的收益分配事宜，通知如下：\n\n"
                f"本期分配方案：\n"
                f"- 分配基准日：{date[:10]}\n"
                f"- 本期应分配收益：{dist_amount_wan}万元\n"
                f"- 年化收益率：{dist_rate_pct}%\n"
                f"- 分配方式：银行转账至委托人指定账户\n"
                f"- 预计到账时间：T+3个工作日\n\n"
                f"请各位确认账户信息是否有变更。\n\n{sender['name']}"
            )
        elif "清算" in topic_name:
            body_text = (
                f"各位好，\n\n"
                f"{trust}已于{date[:10]}到期，现进入清算阶段。\n\n"
                f"清算情况：\n"
                f"- 信托本金：已全额收回\n"
                f"- 累计收益：{liquidation_earnings}万元\n"
                f"- 清算费用：{liquidation_cost}万元\n"
                f"- 预计分配总额：{liquidation_total}万元\n\n"
                f"清算报告详见附件，请审阅。\n\n{sender['name']}"
            )
        else:
            body_text = (
                f"各位同事好，\n\n"
                f"附件为{trust}的{topic_name}，请审阅。\n\n"
                f"主要内容：\n"
                f"- 信托资产运行情况正常\n"
                f"- 融资方按期支付利息\n"
                f"- 抵押物价值稳定\n"
                f"- 合规审查无异常\n\n"
                f"如有疑问请联系信托业务部。\n\n{sender['name']}"
            )

        attachments = []
        if i % 3 == 0:
            fname = f"trust_{topic_key}_{i+1}.docx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_docx(att_path, f"{trust} - {topic_name}", (
                f"华信金融信托有限公司\n{topic_name}\n\n"
                f"信托计划名称：{trust}\n"
                f"编制日期：{date[:10]}\n"
                f"编制部门：信托业务部\n\n"
                f"一、信托计划概况\n\n"
                f"本信托计划由华信金融信托有限公司作为受托人，按照委托人意愿，以受托人名义对信托财产进行管理和处分。"
                f"信托目的为通过专业化的资产管理，实现信托财产的保值增值。\n\n"
                f"二、信托财产管理情况\n\n"
                f"截至报告期末，信托财产总规模为{trust_scale_wan}万元，其中：\n"
                f"- 贷款类资产：{att_loan_assets}万元\n"
                f"- 现金及银行存款：{att_cash_assets}万元\n"
                f"- 应收利息：{att_interest_recv}万元\n\n"
                f"三、风险管理\n\n"
                f"本报告期内，信托项目运行正常，融资方财务状况稳定，抵押物价值充足，"
                f"各项风控措施执行到位。受托人已按照信托合同约定履行管理职责。\n\n"
                f"四、合规情况\n\n"
                f"本信托计划的管理运营符合《信托法》《信托公司管理办法》等法律法规要求，"
                f"已按规定完成信息披露和监管报备。"
            ))
            attachments.append({"filename": fname, "type": "docx"})
        elif i % 3 == 1:
            fname = f"trust_{topic_key}_{i+1}.pptx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_pptx(att_path, f"{trust}", [
                ("Trust Plan Overview",
                 f"Scale: {round(trust_scale_wan / 100, 1)} million RMB\n"
                 f"Duration: {trust_duration_str}\n"
                 f"Expected Return: {trust_rate_pct}% p.a."),
                ("Investment Structure",
                 f"Senior Tranche: {random.randint(60, 80)}%\n"
                 f"Subordinate Tranche: {random.randint(20, 40)}%\n"
                 f"Risk Mitigation: {trust_risk_measure}"),
                ("Risk Assessment",
                 f"Credit Rating: {due_diligence_rating}\n"
                 f"LTV Ratio: {due_diligence_ltv}%\n"
                 f"Debt Service Coverage: {random.uniform(1.2, 2.0):.1f}x"),
            ])
            attachments.append({"filename": fname, "type": "pptx"})
        else:
            fname = f"trust_{topic_key}_{i+1}.xlsx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_xlsx(att_path, "Trust Data", [
                ["Item", "Amount (万元)", "Ratio", "Status", "Due Date"],
                ["Trust Principal", trust_scale_wan, "100%", "Normal", date[:10]],
                ["Accrued Interest", att_interest_recv, f"{trust_rate_pct}%", "Received", date[:10]],
                ["Management Fee", random.randint(50, 300), "0.3%", "Collected", date[:10]],
                ["Custody Fee", random.randint(10, 50), "0.05%", "Collected", date[:10]],
                ["Reserve Fund", random.randint(100, 500), "1%", "Adequate", date[:10]],
            ])
            attachments.append({"filename": fname, "type": "xlsx"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["信托业务", topic_key],
            "attachments": attachments,
        })
    return emails


def gen_fund_emails(start_id: int) -> list[dict]:
    """基金/投资报告 — 15封"""
    emails = []
    fund_names = [
        "华信价值成长混合基金", "华信稳健债券基金", "华信量化对冲基金",
        "华信新能源产业基金", "华信医药健康基金",
    ]
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    report_types = ["净值报告", "投资运作报告", "基金年报", "持仓分析", "业绩归因"]

    for i in range(15):
        eid = f"email_{start_id + i:03d}"
        fund = fund_names[i % len(fund_names)]
        q = quarters[i % 4]
        rtype = report_types[i % 5]
        sender = pick_sender()
        recipients = pick_recipients(sender["email"])
        date = random_date()
        year = date[:4]
        subject = f"[投资] {fund} {year}{q} {rtype}"

        nav = round(random.uniform(0.85, 1.65), 4)
        total_nav = round(nav + random.uniform(0.1, 0.5), 4)
        aum = random.randint(5000, 80000)

        body_text = (
            f"各位好，\n\n"
            f"附件为{fund}{year}年{q}的{rtype}。\n\n"
            f"基金概况：\n"
            f"- 单位净值：{nav}元\n"
            f"- 累计净值：{total_nav}元\n"
            f"- 基金规模：{aum}万元\n"
            f"- 本季度收益率：{random.uniform(-3.0, 8.0):.2f}%\n"
            f"- 同期业绩基准：{random.uniform(-1.0, 5.0):.2f}%\n\n"
            f"市场分析：\n"
            f"本季度{random.choice(['A股市场震荡上行', '债券市场利率下行', '港股市场波动加剧', '大宗商品价格上涨'])}，"
            f"基金通过{random.choice(['增持科技板块', '降低仓位控制回撤', '增配高等级信用债', '调整行业配置'])}应对市场变化。\n\n"
            f"详情见附件。\n\n{sender['name']}\n投资研究部"
        )

        attachments = []
        if i % 3 == 0:
            fname = f"fund_{q}_{year}_{i+1}.xlsx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_xlsx(att_path, f"{q} NAV", [
                ["Date", "Unit NAV", "Cumulative NAV", "Daily Return %", "AUM (万元)"],
                [f"{year}-{i%12+1:02d}-05", nav - 0.02, total_nav - 0.02, f"{random.uniform(-1, 1):.3f}", aum],
                [f"{year}-{i%12+1:02d}-15", nav - 0.01, total_nav - 0.01, f"{random.uniform(-1, 1):.3f}", aum + 100],
                [f"{year}-{i%12+1:02d}-25", nav, total_nav, f"{random.uniform(-1, 1):.3f}", aum + 200],
                ["", "", "", "", ""],
                ["Top Holdings", "Weight %", "Sector", "P/E", "Rating"],
                [random.choice(["Kweichow Moutai", "CATL", "BYD"]), f"{random.randint(3,8)}%", "Consumer", random.randint(15, 40), "Buy"],
                [random.choice(["China Merchants Bank", "Ping An", "ICBC"]), f"{random.randint(3,8)}%", "Finance", random.randint(5, 15), "Hold"],
                [random.choice(["Longi Green", "Sungrow", "Ganfeng Lithium"]), f"{random.randint(2,6)}%", "New Energy", random.randint(20, 60), "Buy"],
            ])
            attachments.append({"filename": fname, "type": "xlsx"})
        elif i % 3 == 1:
            fname = f"fund_{q}_{year}_{i+1}.pdf"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_pdf(att_path, f"{fund} {q} Report", (
                f"Fund Performance Report\n\n"
                f"Fund Name: {fund}\n"
                f"Report Period: {year} {q}\n"
                f"NAV per Unit: {nav}\n"
                f"AUM: RMB {aum}0,000\n\n"
                f"Asset Allocation:\n"
                f"- Equities: {random.randint(50, 80)}%\n"
                f"- Fixed Income: {random.randint(10, 30)}%\n"
                f"- Cash: {random.randint(5, 15)}%\n"
                f"- Others: {random.randint(0, 10)}%\n\n"
                f"Performance Attribution:\n"
                f"- Sector allocation: +{random.uniform(0.5, 3.0):.2f}%\n"
                f"- Stock selection: {random.uniform(-1.0, 2.0):+.2f}%\n"
                f"- Timing: {random.uniform(-0.5, 1.0):+.2f}%\n\n"
                f"Risk Metrics:\n"
                f"- Sharpe Ratio: {random.uniform(0.5, 2.0):.2f}\n"
                f"- Max Drawdown: -{random.uniform(3, 15):.2f}%\n"
                f"- Volatility: {random.uniform(8, 25):.2f}%"
            ))
            attachments.append({"filename": fname, "type": "pdf"})
        else:
            fname = f"fund_holdings_{q}_{year}_{i+1}.csv"
            att_path = str(ATTACHMENTS_DIR / fname)
            stocks = [
                ("600519", "贵州茅台", "消费"), ("300750", "宁德时代", "新能源"),
                ("601318", "中国平安", "金融"), ("002594", "比亚迪", "汽车"),
                ("600036", "招商银行", "金融"), ("601012", "隆基绿能", "新能源"),
                ("000858", "五粮液", "消费"), ("002415", "海康威视", "科技"),
            ]
            selected = random.sample(stocks, min(6, len(stocks)))
            make_csv_file(att_path,
                ["Stock Code", "Stock Name", "Sector", "Weight %", "Market Value (万元)", "P/E Ratio"],
                [[s[0], s[1], s[2], f"{random.uniform(2, 8):.1f}", str(random.randint(200, 3000)), str(random.randint(10, 50))]
                 for s in selected]
            )
            attachments.append({"filename": fname, "type": "csv"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [{"name": PEOPLE[6]["name"], "email": PEOPLE[6]["email"]}],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["投资", fund],
            "attachments": attachments,
        })
    return emails


def gen_compliance_emails(start_id: int) -> list[dict]:
    """合规/监管 — 12封"""
    emails = []
    topics = [
        ("银保监会信托监管评级结果通知", "regulatory_rating"),
        ("反洗钱自查报告", "aml_review"),
        ("关联交易合规审查", "related_party"),
        ("信托产品信息披露检查", "disclosure"),
        ("投资者适当性管理自查", "suitability"),
        ("资管新规整改进度报告", "asset_mgmt_reform"),
        ("净资本管理月度报告", "net_capital"),
        ("流动性风险监测报告", "liquidity"),
        ("操作风险事件排查", "op_risk"),
        ("数据治理合规报告", "data_governance"),
        ("监管处罚案例学习", "regulatory_case"),
        ("内控制度修订通知", "internal_control"),
    ]

    for i in range(12):
        eid = f"email_{start_id + i:03d}"
        topic_name, topic_key = topics[i]
        sender = PEOPLE[2] if i % 2 == 0 else PEOPLE[7]  # 风控合规部 or 合规审计部
        recipients = pick_recipients(sender["email"], count=random.randint(3, 6))
        date = random_date()
        subject = f"[合规] {topic_name}"

        body_text = (
            f"各位同事，\n\n"
            f"{'现将' + topic_name + '发送如下' if i % 2 == 0 else '关于' + topic_name + '，通知如下'}：\n\n"
        )

        if "反洗钱" in topic_name:
            body_text += (
                f"根据《反洗钱法》及银保监会相关规定，我司已完成{date[:4]}年度反洗钱自查工作：\n\n"
                f"1. 客户身份识别：已对全部{random.randint(500, 2000)}名客户完成身份核实\n"
                f"2. 大额交易报告：本期报送{random.randint(20, 100)}笔\n"
                f"3. 可疑交易报告：本期报送{random.randint(5, 30)}笔\n"
                f"4. 客户风险等级分类：高风险客户{random.randint(10, 50)}户\n\n"
                f"详见附件报告。\n"
            )
        elif "净资本" in topic_name:
            net_capital = random.randint(20, 80)
            body_text += (
                f"截至{date[:10]}，我司净资本情况如下：\n\n"
                f"- 净资本：{net_capital}亿元\n"
                f"- 各项业务风险资本之和：{random.randint(10, 40)}亿元\n"
                f"- 净资本/各项业务风险资本之和：{random.randint(150, 300)}%（监管要求≥100%）\n"
                f"- 净资本/净资产：{random.randint(70, 95)}%（监管要求≥40%）\n\n"
                f"各项指标均满足监管要求。\n"
            )
        elif "资管新规" in topic_name:
            body_text += (
                f"按照资管新规过渡期安排，我司整改进度如下：\n\n"
                f"1. 非标资产压降：已压降至{random.randint(100, 500)}亿，完成目标的{random.randint(70, 95)}%\n"
                f"2. 多层嵌套清理：{random.randint(85, 100)}%已完成\n"
                f"3. 资金池业务：已全部清理\n"
                f"4. 估值方法切换：{random.randint(90, 100)}%产品已完成\n\n"
            )
        else:
            body_text += (
                f"请各部门认真学习并落实相关要求。\n"
                f"如有违规事项，请立即上报合规部门。\n\n"
                f"合规要点：\n"
                f"- 严格执行投资者适当性管理要求\n"
                f"- 确保信息披露的真实、准确、完整、及时\n"
                f"- 加强关联交易管控，防范利益输送\n"
            )

        body_text += f"\n{sender['name']}\n{sender['dept']}"

        attachments = []
        if i % 2 == 0:
            fname = f"compliance_{topic_key}_{i+1}.docx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_docx(att_path, topic_name, (
                f"华信金融信托有限公司\n{topic_name}\n\n"
                f"报告期间：{date[:7]}\n编制部门：{sender['dept']}\n\n"
                f"一、总体情况\n\n"
                f"报告期内，公司各项业务运营合规，内部控制有效执行。"
                f"根据监管要求和公司内控制度，对各业务条线进行了全面合规检查。\n\n"
                f"二、检查发现\n\n"
                f"本次检查共涉及{random.randint(10, 50)}项检查要点，"
                f"发现{random.randint(0, 5)}项需改进事项，均为非重大问题。\n\n"
                f"三、整改措施\n\n"
                f"1. 完善相关业务流程文档\n"
                f"2. 加强员工合规培训\n"
                f"3. 优化系统控制节点\n\n"
                f"四、后续计划\n\n"
                f"将持续跟踪整改进度，确保在规定期限内完成全部整改工作。"
            ))
            attachments.append({"filename": fname, "type": "docx"})
        else:
            fname = f"compliance_{topic_key}_{i+1}.txt"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_txt(att_path, (
                f"== {topic_name} ==\n"
                f"日期: {date[:10]}\n"
                f"编制: {sender['name']} ({sender['dept']})\n\n"
                f"摘要:\n"
                f"本报告对公司{topic_name.replace('报告', '').replace('通知', '')}相关事项进行了全面梳理。\n\n"
                f"关键指标:\n"
                f"- 合规检查覆盖率: {random.randint(95, 100)}%\n"
                f"- 问题整改完成率: {random.randint(85, 100)}%\n"
                f"- 员工合规培训完成率: {random.randint(90, 100)}%\n"
                f"- 可疑交易监测有效率: {random.randint(95, 100)}%\n\n"
                f"结论:\n"
                f"公司合规管理体系运行有效，各项监管指标满足要求。"
            ))
            attachments.append({"filename": fname, "type": "txt"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["合规", topic_key],
            "attachments": attachments,
        })
    return emails


def gen_risk_emails(start_id: int) -> list[dict]:
    """风控报告 — 10封"""
    emails = []
    risk_topics = [
        "信用风险监测周报", "市场风险日报", "流动性风险监测报告",
        "集中度风险分析", "信托项目逾期预警", "压力测试报告",
        "风险限额使用情况", "大额风险暴露报告", "风险事件应急预案",
        "风险偏好声明书",
    ]

    for i in range(10):
        eid = f"email_{start_id + i:03d}"
        topic = risk_topics[i]
        sender = PEOPLE[2]  # 风控合规部
        recipients = pick_recipients(sender["email"], count=random.randint(2, 5))
        date = random_date()
        subject = f"[风控] {topic} - {date[:10]}"

        body_text = (
            f"各位领导好，\n\n"
            f"附件为{topic}，截至{date[:10]}。\n\n"
        )

        if "信用风险" in topic:
            body_text += (
                f"信用风险概况：\n"
                f"- 在管信托项目：{random.randint(50, 200)}个\n"
                f"- 正常类：{random.randint(80, 95)}%\n"
                f"- 关注类：{random.randint(3, 10)}%\n"
                f"- 不良类：{random.randint(1, 5)}%\n"
                f"- 加权平均信用评级：{random.choice(['AA', 'AA-', 'A+'])}\n"
            )
        elif "逾期" in topic:
            body_text += (
                f"逾期预警信息：\n"
                f"- 本周新增逾期项目：{random.randint(0, 3)}个\n"
                f"- 逾期金额：{random.randint(0, 5000)}万元\n"
                f"- 已采取催收措施：{random.choice(['电话催收', '发送催收函', '启动法律程序'])}\n"
                f"- 预计回收情况：{random.choice(['预计30日内全额回收', '已制定分期还款方案', '需要进一步评估'])}\n"
            )
        elif "压力测试" in topic:
            body_text += (
                f"压力测试结果：\n"
                f"- 轻度压力情景：资本充足率{random.randint(130, 180)}%，满足监管要求\n"
                f"- 中度压力情景：资本充足率{random.randint(100, 130)}%，满足监管要求\n"
                f"- 重度压力情景：资本充足率{random.randint(80, 100)}%，{random.choice(['仍满足监管要求', '接近监管红线，需关注'])}\n"
            )
        else:
            body_text += (
                f"报告要点：\n"
                f"- 各项风险指标在可控范围内\n"
                f"- 未发现重大风险隐患\n"
                f"- 风险限额使用率：{random.randint(40, 80)}%\n"
            )

        body_text += f"\n详见附件。\n\n{sender['name']}\n风控合规部"

        attachments = []
        if i % 2 == 0:
            fname = f"risk_{topic.replace('/', '_')}_{i+1}.xlsx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_xlsx(att_path, "Risk Monitor", [
                ["Project", "Scale (万元)", "Risk Level", "Rating", "Maturity", "Status"],
                *[
                    [f"Trust-{random.randint(100,999)}", random.randint(5000, 50000),
                     random.choice(["Low", "Medium", "Low", "Low"]),
                     random.choice(["AA+", "AA", "AA-", "A+"]),
                     f"{date[:4]}-{random.randint(1,12):02d}-{random.randint(1,28):02d}",
                     random.choice(["Normal", "Normal", "Normal", "Watch"])]
                    for _ in range(8)
                ],
            ])
            attachments.append({"filename": fname, "type": "xlsx"})
        else:
            fname = f"risk_{topic.replace('/', '_')}_{i+1}.pdf"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_pdf(att_path, f"Risk Report - {date[:10]}", (
                f"Huaxin Trust Risk Management Report\n\n"
                f"Report: {topic}\n"
                f"Date: {date[:10]}\n\n"
                f"1. Risk Overview\n"
                f"Total AUM under monitoring: RMB {random.randint(100, 800)} billion\n"
                f"Number of active projects: {random.randint(50, 200)}\n\n"
                f"2. Key Risk Indicators\n"
                f"- NPL Ratio: {random.uniform(0.5, 3.0):.2f}%\n"
                f"- Provision Coverage: {random.randint(150, 300)}%\n"
                f"- Concentration Ratio (Top 10): {random.randint(20, 45)}%\n"
                f"- Liquidity Coverage Ratio: {random.randint(120, 200)}%\n\n"
                f"3. Risk Events\n"
                f"No significant risk events during the reporting period.\n\n"
                f"4. Recommendations\n"
                f"- Continue monitoring macro-economic indicators\n"
                f"- Enhance post-investment management for real estate projects\n"
                f"- Review counterparty credit limits quarterly"
            ))
            attachments.append({"filename": fname, "type": "pdf"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["风控", topic],
            "attachments": attachments,
        })
    return emails


def gen_legal_emails(start_id: int) -> list[dict]:
    """法务/合同 — 12封"""
    emails = []
    topics = [
        ("信托合同审批", "trust_contract"),
        ("投资顾问协议", "advisor_agreement"),
        ("托管协议定稿", "custody_agreement"),
        ("担保合同审查", "guarantee_contract"),
        ("股权质押合同", "pledge_contract"),
        ("保密协议(NDA)", "nda"),
        ("资产转让协议", "asset_transfer"),
        ("基金合同修订", "fund_amendment"),
        ("合规法律意见书", "legal_opinion"),
        ("诉讼案件进展通报", "litigation"),
        ("监管问询回复函", "regulatory_response"),
        ("知识产权许可协议", "ip_license"),
    ]

    for i in range(12):
        eid = f"email_{start_id + i:03d}"
        topic_name, topic_key = topics[i]
        sender = PEOPLE[3]  # 法务部 赵雪梅
        recipients = pick_recipients(sender["email"])
        date = random_date()
        contract_no = f"HX-{date[:4]}-{random.randint(1000, 9999)}"
        subject = f"[法务] {topic_name} - {contract_no}"

        body_text = (
            f"各位好，\n\n"
            f"附件为{topic_name}（编号：{contract_no}），请审阅。\n\n"
        )

        if "信托合同" in topic_name:
            body_text += (
                f"合同要点：\n"
                f"- 信托类型：{random.choice(['集合资金信托', '单一资金信托', '财产权信托'])}\n"
                f"- 信托规模：{random.randint(5000, 50000)}万元\n"
                f"- 信托期限：{random.choice(['12', '18', '24', '36'])}个月\n"
                f"- 受益人类型：{random.choice(['单一受益人', '多个受益人（同一顺序）', '优先/劣后受益人'])}\n"
                f"- 信托报酬率：{random.uniform(0.3, 1.5):.2f}%/年\n"
            )
        elif "诉讼" in topic_name:
            body_text += (
                f"案件进展：\n"
                f"- 案由：{random.choice(['金融借款合同纠纷', '信托纠纷', '担保合同纠纷', '票据纠纷'])}\n"
                f"- 涉案金额：{random.randint(500, 10000)}万元\n"
                f"- 当前阶段：{random.choice(['已立案', '举证阶段', '一审开庭', '等待判决', '执行阶段'])}\n"
                f"- 预计结果：{random.choice(['胜诉概率较大', '正在积极调解', '已达成和解意向'])}\n"
            )
        elif "监管问询" in topic_name:
            body_text += (
                f"问询事项：\n"
                f"- 问询机关：{random.choice(['银保监会', '地方银保监局', '中国信托登记有限责任公司'])}\n"
                f"- 问询内容：{random.choice(['个别信托项目运行情况', '关联交易情况说明', '净资本计算方法'])}\n"
                f"- 回复期限：收函后{random.randint(5, 15)}个工作日\n"
                f"- 需要各部门配合提供材料，具体见附件。\n"
            )
        else:
            body_text += (
                f"合同金额：{random.randint(100, 50000)}万元\n"
                f"签署方：华信金融信托有限公司 与 {random.choice(['恒丰投资有限公司', '龙泰财富管理有限公司', '中原资本集团'])}\n"
                f"请在{random.randint(3, 7)}个工作日内完成审阅。\n"
            )

        body_text += f"\n{sender['name']}\n法务部"

        attachments = []
        if i % 2 == 0:
            fname = f"legal_{topic_key}_{contract_no}.pdf"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_pdf(att_path, f"Legal Document - {contract_no}", (
                f"Huaxin Trust Co., Ltd.\n"
                f"Document: {topic_name}\n"
                f"Reference: {contract_no}\n"
                f"Date: {date[:10]}\n\n"
                f"ARTICLE 1 - PARTIES\n"
                f"Trustee: Huaxin Financial Trust Co., Ltd.\n"
                f"Settlor/Beneficiary: As specified in the Trust Deed\n\n"
                f"ARTICLE 2 - PURPOSE\n"
                f"This agreement establishes the terms for trust asset management.\n\n"
                f"ARTICLE 3 - TRUST PROPERTY\n"
                f"The trust property shall be managed in accordance with the trust deed.\n\n"
                f"ARTICLE 4 - RIGHTS AND OBLIGATIONS\n"
                f"The trustee shall exercise fiduciary duty with due care and diligence.\n\n"
                f"ARTICLE 5 - FEES\n"
                f"Management fee, custody fee, and performance fee as specified.\n\n"
                f"ARTICLE 6 - TERMINATION\n"
                f"This agreement terminates upon expiry or early termination conditions."
            ))
            attachments.append({"filename": fname, "type": "pdf"})
        else:
            fname = f"legal_{topic_key}_{contract_no}.docx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_docx(att_path, f"{topic_name} - {contract_no}", (
                f"合同编号：{contract_no}\n签订日期：{date[:10]}\n\n"
                f"甲方（受托人）：华信金融信托有限公司\n"
                f"乙方：[签约对方]\n\n"
                f"第一条 合同目的\n\n"
                f"本合同就{topic_name}相关事项达成如下协议，明确双方权利义务。\n\n"
                f"第二条 信托财产\n\n"
                f"信托财产的范围、管理方式、处分规则按照信托合同约定执行。"
                f"受托人应以受益人利益最大化为原则管理信托财产。\n\n"
                f"第三条 信息披露\n\n"
                f"受托人应按照法律法规和监管要求，及时、准确、完整地向委托人和受益人披露信息。\n\n"
                f"第四条 费用\n\n"
                f"信托报酬：{random.uniform(0.3, 1.5):.2f}%/年\n"
                f"托管费：{random.uniform(0.03, 0.1):.2f}%/年\n"
                f"其他费用按实际发生额收取。"
            ))
            attachments.append({"filename": fname, "type": "docx"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [{"name": PEOPLE[7]["name"], "email": PEOPLE[7]["email"]}],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["法务", topic_key],
            "attachments": attachments,
        })
    return emails


def gen_operations_emails(start_id: int) -> list[dict]:
    """运营管理/内部沟通 — 12封"""
    emails = []
    topics = [
        ("信托产品到期兑付通知", "redemption"),
        ("新员工入职培训-信托基础知识", "training"),
        ("IT系统升级通知-信托核算系统", "it_upgrade"),
        ("客户投诉处理流程优化", "complaint"),
        ("季度经营分析会纪要", "quarterly_review"),
        ("信托登记系统操作规范", "registration_ops"),
        ("年度合规培训安排", "compliance_training"),
        ("办公场所安全检查", "safety"),
        ("员工持证情况统计", "certification"),
        ("部门年度预算编制", "budget"),
        ("信托业务系统故障应急处理", "incident"),
        ("跨部门协作流程规范", "workflow"),
    ]

    for i in range(12):
        eid = f"email_{start_id + i:03d}"
        topic_name, topic_key = topics[i]
        sender = pick_sender()
        recipients = pick_recipients(sender["email"], count=random.randint(2, 6))
        date = random_date()
        subject = f"[运营] {topic_name}"

        body_text = (
            f"各位同事，\n\n"
        )

        if "兑付" in topic_name:
            body_text += (
                f"以下信托产品即将到期，请做好兑付准备：\n\n"
                f"1. 华信·稳健增长1号 - 到期日：{date[:10]}，规模：{random.randint(5000, 20000)}万元\n"
                f"2. 华信·城市发展3号 - 到期日：{date[:8]}28，规模：{random.randint(3000, 15000)}万元\n\n"
                f"请运营部核实资金回款情况，财务部准备兑付资金，客服部做好客户通知工作。\n"
            )
        elif "培训" in topic_name:
            body_text += (
                f"本次培训安排如下：\n\n"
                f"主题：{random.choice(['信托法律法规', '反洗钱实务', '信托产品设计', '风险管理框架', '投资者适当性管理'])}\n"
                f"时间：{date[:10]} 14:00-17:00\n"
                f"地点：{random.choice(['公司12楼培训室', '线上腾讯会议', '公司15楼大会议室'])}\n"
                f"讲师：{random.choice(['外部专家', '合规部总经理', '业务部资深经理'])}\n\n"
                f"请相关人员准时参加，签到考勤。\n"
            )
        elif "经营分析" in topic_name:
            body_text += (
                f"季度经营分析会主要议题：\n\n"
                f"1. 信托资产管理规模：{random.randint(500, 2000)}亿元，同比{random.choice(['增长', '下降'])}{random.randint(3, 15)}%\n"
                f"2. 信托报酬收入：{random.randint(2, 15)}亿元\n"
                f"3. 新增项目：{random.randint(10, 50)}个，规模{random.randint(50, 300)}亿元\n"
                f"4. 到期兑付：{random.randint(20, 80)}个项目正常兑付\n"
                f"5. 不良信托资产处置进展\n\n"
                f"详细纪要见附件。\n"
            )
        elif "持证" in topic_name:
            body_text += (
                f"请各部门配合统计员工持证情况：\n\n"
                f"需统计的证书包括：\n"
                f"- 基金从业资格证\n"
                f"- 证券从业资格证\n"
                f"- 银行从业资格证\n"
                f"- CFA/CPA/FRM\n"
                f"- 法律职业资格证\n\n"
                f"请于{date[:8]}28前将统计表反馈至人事部。\n"
            )
        else:
            body_text += (
                f"关于{topic_name}，请各部门配合执行。\n"
                f"详细内容和要求见附件。\n"
            )

        body_text += f"\n{sender['name']}"

        attachments = []
        if "截图" in topic_name or "系统" in topic_name:
            if "故障" in topic_name:
                fname = f"ops_{topic_key}_{i+1}.png"
                att_path = str(ATTACHMENTS_DIR / fname)
                make_png(att_path, f"System Alert: {topic_name}\n\nIncident Time: {date}\nSeverity: High\nStatus: Resolved")
                attachments.append({"filename": fname, "type": "png"})
            else:
                fname = f"ops_{topic_key}_{i+1}.jpg"
                att_path = str(ATTACHMENTS_DIR / fname)
                make_jpg(att_path, f"System Screenshot\n{topic_name}\n{date[:10]}")
                attachments.append({"filename": fname, "type": "jpg"})
        elif i % 2 == 0:
            fname = f"ops_{topic_key}_{i+1}.docx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_docx(att_path, topic_name, (
                f"华信金融信托有限公司\n{topic_name}\n\n"
                f"日期：{date[:10]}\n编制部门：{sender['dept']}\n\n"
                f"一、背景\n\n"
                f"为规范公司运营管理流程，提升工作效率，特制定本方案。\n\n"
                f"二、具体内容\n\n"
                f"1. 加强信托产品全生命周期管理\n"
                f"2. 优化内部审批流程，提高效率\n"
                f"3. 完善信息系统建设，提升数字化水平\n"
                f"4. 加强人才培养和团队建设\n\n"
                f"三、执行要求\n\n"
                f"请各部门严格按照方案要求执行，确保各项工作落实到位。"
            ))
            attachments.append({"filename": fname, "type": "docx"})
        else:
            fname = f"ops_{topic_key}_{i+1}.txt"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_txt(att_path, (
                f"== {topic_name} ==\n"
                f"日期: {date[:10]}\n\n"
                f"要点:\n"
                f"- 严格执行各项管理制度\n"
                f"- 加强部门间沟通协作\n"
                f"- 按时完成各项工作任务\n"
                f"- 及时反馈执行中的问题\n\n"
                f"联系人: {sender['name']} ({sender['dept']})\n"
                f"电话: 010-8888-{random.randint(1000,9999)}"
            ))
            attachments.append({"filename": fname, "type": "txt"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["运营", topic_key],
            "attachments": attachments,
        })
    return emails


def gen_wealth_emails(start_id: int) -> list[dict]:
    """财富管理/客户服务 — 10封"""
    emails = []
    topics = [
        "高净值客户资产配置方案", "家族信托设立咨询", "慈善信托方案设计",
        "保险金信托合作方案", "客户KYC信息更新通知", "投资者教育活动方案",
        "客户满意度调查报告", "理财经理考核标准", "客户资产配置季度回顾",
        "高端客户服务标准手册",
    ]

    for i in range(10):
        eid = f"email_{start_id + i:03d}"
        topic = topics[i]
        sender = pick_sender()
        recipients = pick_recipients(sender["email"])
        date = random_date()
        subject = f"[财富管理] {topic}"

        body_text = (
            f"各位好，\n\n"
        )

        if "家族信托" in topic:
            body_text += (
                f"关于家族信托设立咨询事项，整理如下：\n\n"
                f"家族信托基本方案：\n"
                f"- 信托规模：不低于{random.choice(['1000万', '3000万', '5000万', '1亿'])}元\n"
                f"- 信托期限：{random.choice(['永续', '30年', '50年'])}\n"
                f"- 受益人安排：可设定多层受益人，包括配偶、子女、孙辈\n"
                f"- 分配规则：可按月/按季/按年分配，或设定条件分配（如教育、婚姻、创业等）\n"
                f"- 资产类型：现金、股权、不动产、保险金等均可装入信托\n\n"
                f"家族信托的核心优势：\n"
                f"1. 资产隔离保护 — 信托财产独立于委托人、受托人、受益人的固有财产\n"
                f"2. 财富传承规划 — 按照委托人意愿实现跨代际财富传承\n"
                f"3. 税务筹划 — 合理的信托架构可实现税务优化\n"
                f"4. 隐私保护 — 信托信息不公开\n\n"
                f"详细方案见附件。\n"
            )
        elif "慈善信托" in topic:
            body_text += (
                f"慈善信托方案要点：\n\n"
                f"- 信托目的：{random.choice(['教育扶贫', '医疗救助', '环境保护', '科技创新奖励'])}\n"
                f"- 设立规模：{random.randint(100, 5000)}万元\n"
                f"- 信托期限：永续\n"
                f"- 受益人：不特定社会公众\n"
                f"- 监察人：{random.choice(['民政局指定', '委托人推荐的第三方机构'])}\n\n"
                f"根据《慈善法》和银保监会相关规定，慈善信托需在民政部门备案。\n"
            )
        elif "资产配置" in topic:
            body_text += (
                f"本季度资产配置建议：\n\n"
                f"保守型客户：\n"
                f"- 固收类信托：{random.randint(40, 60)}%\n"
                f"- 债券基金：{random.randint(20, 30)}%\n"
                f"- 现金管理：{random.randint(10, 20)}%\n\n"
                f"平衡型客户：\n"
                f"- 固收类信托：{random.randint(30, 40)}%\n"
                f"- 混合基金：{random.randint(20, 30)}%\n"
                f"- 股权投资：{random.randint(10, 20)}%\n"
                f"- 现金管理：{random.randint(10, 15)}%\n\n"
                f"进取型客户：\n"
                f"- 股权投资信托：{random.randint(30, 40)}%\n"
                f"- 权益类基金：{random.randint(20, 30)}%\n"
                f"- 另类投资：{random.randint(10, 20)}%\n"
                f"- 固收类：{random.randint(10, 15)}%\n"
            )
        elif "KYC" in topic:
            body_text += (
                f"根据反洗钱法规要求，以下客户的KYC信息需要更新：\n\n"
                f"- 需更新客户数量：{random.randint(50, 200)}户\n"
                f"- 更新内容：身份证件、职业信息、资产证明、风险偏好\n"
                f"- 截止日期：{date[:8]}28\n\n"
                f"请理财经理尽快联系客户完成信息更新。\n"
            )
        else:
            body_text += f"详细内容见附件。\n"

        body_text += f"\n{sender['name']}\n{sender['dept']}"

        attachments = []
        if i % 2 == 0:
            fname = f"wealth_{topic.replace('/', '_')}_{i+1}.pdf"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_pdf(att_path, f"Wealth Management - {topic}", (
                f"Huaxin Trust Wealth Management\n\n"
                f"Document: {topic}\n"
                f"Date: {date[:10]}\n\n"
                f"1. Overview\n"
                f"Huaxin Trust provides comprehensive wealth management services\n"
                f"for high-net-worth individuals and families.\n\n"
                f"2. Service Offerings\n"
                f"- Family Trust: Asset protection and succession planning\n"
                f"- Charitable Trust: Philanthropic giving with tax benefits\n"
                f"- Insurance Trust: Life insurance integrated with trust\n"
                f"- Investment Advisory: Customized portfolio management\n\n"
                f"3. Client Segmentation\n"
                f"- Ultra HNW (>50M RMB): Dedicated family office service\n"
                f"- HNW (10-50M RMB): Senior wealth advisor\n"
                f"- Affluent (3-10M RMB): Personal wealth advisor\n\n"
                f"4. Compliance\n"
                f"All services comply with Trust Law and CBIRC regulations."
            ))
            attachments.append({"filename": fname, "type": "pdf"})
        else:
            fname = f"wealth_{topic.replace('/', '_')}_{i+1}.xlsx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_xlsx(att_path, "Client Data", [
                ["Client Type", "Count", "AUM (亿元)", "Avg AUM (万元)", "Products Held"],
                ["Ultra HNW", random.randint(20, 80), random.randint(50, 200), random.randint(5000, 20000), random.randint(3, 8)],
                ["HNW", random.randint(100, 500), random.randint(100, 500), random.randint(1000, 5000), random.randint(2, 5)],
                ["Affluent", random.randint(500, 2000), random.randint(200, 800), random.randint(300, 1000), random.randint(1, 3)],
                ["", "", "", "", ""],
                ["Product Type", "AUM (亿元)", "Client Count", "Avg Return", "Maturity"],
                ["Fixed Income Trust", random.randint(100, 400), random.randint(200, 800), f"{random.uniform(5, 7):.1f}%", "12-24M"],
                ["Equity Trust", random.randint(30, 100), random.randint(50, 200), f"{random.uniform(8, 15):.1f}%", "24-36M"],
                ["Family Trust", random.randint(50, 200), random.randint(20, 80), "N/A", "Perpetual"],
            ])
            attachments.append({"filename": fname, "type": "xlsx"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["财富管理", topic],
            "attachments": attachments,
        })
    return emails


def gen_research_emails(start_id: int) -> list[dict]:
    """投研/市场分析 — 14封"""
    emails = []
    topics = [
        ("宏观经济月度研判", "macro"),
        ("房地产行业分析报告", "real_estate"),
        ("地方政府融资平台分析", "lgfv"),
        ("信托行业发展趋势报告", "industry_trend"),
        ("利率走势分析与预测", "interest_rate"),
        ("股票市场策略周报", "equity_strategy"),
        ("债券市场分析报告", "bond_market"),
        ("另类投资机会研究", "alternatives"),
        ("ESG投资研究报告", "esg"),
        ("科技行业投资机会", "tech_sector"),
        ("消费行业深度研究", "consumer"),
        ("新能源产业链分析", "new_energy"),
        ("医药行业政策解读", "pharma_policy"),
        ("海外市场配置建议", "overseas"),
    ]

    for i in range(14):
        eid = f"email_{start_id + i:03d}"
        topic_name, topic_key = topics[i]
        sender = pick_sender()
        recipients = pick_recipients(sender["email"], count=random.randint(3, 8))
        date = random_date()
        subject = f"[投研] {topic_name} - {date[:7]}"

        # 预先生成宏观数据共享变量
        macro_gdp = round(random.uniform(4.5, 6.5), 1)
        macro_gdp_prev = round(random.uniform(4.0, 6.0), 1)
        macro_cpi = round(random.uniform(0.5, 3.0), 1)
        macro_pmi = round(random.uniform(49.0, 52.0), 1)
        macro_pmi_desc = random.choice(['处于扩张区间', '接近荣枯线'])
        macro_social_finance = round(random.uniform(9.0, 12.0), 1)

        body_text = (
            f"各位同事好，\n\n"
            f"附件为最新的{topic_name}，核心观点如下：\n\n"
        )

        if "宏观" in topic_name:
            body_text += (
                f"1. GDP增速预测：{macro_gdp}%（前值{macro_gdp_prev}%）\n"
                f"2. CPI同比：{macro_cpi}%，通胀温和\n"
                f"3. PMI：{macro_pmi}，{macro_pmi_desc}\n"
                f"4. 社融增速：{macro_social_finance}%\n"
                f"5. 政策展望：{random.choice(['货币政策保持稳健偏松', '财政政策积极发力', '预计将出台定向降准'])}\n"
            )
        elif "房地产" in topic_name:
            body_text += (
                f"1. 商品房销售面积同比{random.choice(['增长', '下降'])}{random.randint(3, 20)}%\n"
                f"2. 土地出让收入同比{random.choice(['增长', '下降'])}{random.randint(5, 30)}%\n"
                f"3. 房企融资环境：{random.choice(['边际改善', '仍然偏紧', '分化明显'])}\n"
                f"4. 政策面：{random.choice(['限购政策有所放松', '保交楼政策持续推进', 'LPR下调利好需求端'])}\n"
                f"5. 投资建议：{random.choice(['关注优质房企的融资类信托', '审慎介入，控制集中度', '优选一二线城市项目'])}\n"
            )
        elif "地方政府" in topic_name or "lgfv" in topic_key:
            body_text += (
                f"1. 城投债发行规模：{random.randint(3000, 8000)}亿元\n"
                f"2. 城投平台信用利差：{random.randint(50, 200)}BP\n"
                f"3. 重点关注区域：{random.choice(['江浙地区城投整体稳健', '西南地区需警惕尾部风险', '山东区域分化加剧'])}\n"
                f"4. 隐性债务化解进展：{random.choice(['化债方案稳步推进', '特殊再融资债发行加速', '部分区域仍有压力'])}\n"
                f"5. 投资建议：优选经济财力强的区域，控制平台层级\n"
            )
        elif "利率" in topic_name:
            body_text += (
                f"1. 10年期国债收益率：{random.uniform(2.5, 3.5):.2f}%\n"
                f"2. 1年期MLF利率：{random.uniform(2.5, 3.0):.2f}%\n"
                f"3. DR007：{random.uniform(1.5, 2.5):.2f}%\n"
                f"4. 预测：{random.choice(['利率中枢下移趋势延续', '短期利率或有小幅反弹', '长端利率已处历史低位，下行空间有限'])}\n"
            )
        else:
            body_text += (
                f"1. 行业增速：{random.randint(5, 30)}%\n"
                f"2. 市场规模：{random.randint(1000, 50000)}亿元\n"
                f"3. 投资机会：{random.choice(['行业处于高速发展期', '估值处于历史低位', '政策红利释放中'])}\n"
                f"4. 风险提示：{random.choice(['政策不确定性', '技术路径变化', '竞争格局重塑', '估值偏高'])}\n"
            )

        body_text += f"\n请参考附件了解详情。\n\n{sender['name']}\n投资研究部"

        attachments = []
        if i % 3 == 0:
            fname = f"research_{topic_key}_{i+1}.pdf"
            att_path = str(ATTACHMENTS_DIR / fname)
            if "宏观" in topic_name:
                macro_content = (
                    f"Huaxin Trust Investment Research\n\n"
                    f"Report: {topic_name}\n"
                    f"Date: {date[:10]}\n"
                    f"Analyst: {sender['name']}\n\n"
                    f"Macroeconomic Key Indicators\n\n"
                    f"GDP Growth Forecast: {macro_gdp}% (prev: {macro_gdp_prev}%)\n"
                    f"CPI YoY: {macro_cpi}% (mild inflation)\n"
                    f"PMI: {macro_pmi} ({macro_pmi_desc})\n"
                    f"Social Financing Growth: {macro_social_finance}%\n\n"
                    f"Investment Recommendations:\n"
                    f"- Maintain {random.choice(['overweight', 'neutral', 'underweight'])} position\n"
                    f"- Target allocation: {random.randint(5, 20)}% of portfolio\n"
                    f"- Risk rating: {random.choice(['Low', 'Medium', 'Medium-High'])}"
                )
            else:
                macro_content = (
                    f"Huaxin Trust Investment Research\n\n"
                    f"Report: {topic_name}\n"
                    f"Date: {date[:10]}\n"
                    f"Analyst: {sender['name']}\n\n"
                    f"Executive Summary\n"
                    f"This report provides analysis on {topic_name}.\n\n"
                    f"Key Findings:\n"
                    f"- Market conditions remain {random.choice(['favorable', 'challenging', 'mixed'])}\n"
                    f"- Sector rotation suggests {random.choice(['growth to value shift', 'defensive positioning', 'cyclical recovery'])}\n"
                    f"- Risk-adjusted returns are {random.choice(['above average', 'in line with expectations', 'below historical norms'])}\n\n"
                    f"Investment Recommendations:\n"
                    f"- Maintain {random.choice(['overweight', 'neutral', 'underweight'])} position\n"
                    f"- Target allocation: {random.randint(5, 20)}% of portfolio\n"
                    f"- Risk rating: {random.choice(['Low', 'Medium', 'Medium-High'])}"
                )
            make_pdf(att_path, f"Research Report - {topic_name}", macro_content)
            attachments.append({"filename": fname, "type": "pdf"})
        elif i % 3 == 1:
            fname = f"research_{topic_key}_{i+1}.docx"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_docx(att_path, f"投研报告 - {topic_name}", (
                f"华信金融信托有限公司\n投资研究部\n\n"
                f"报告主题：{topic_name}\n"
                f"报告日期：{date[:10]}\n"
                f"分析师：{sender['name']}\n\n"
                f"一、市场回顾\n\n"
                f"本报告期内，市场整体呈现{random.choice(['震荡上行', '窄幅整理', '先扬后抑', '稳步回升'])}态势。"
                f"主要受{random.choice(['货币政策宽松', '经济数据改善', '外资持续流入', '政策预期提振'])}驱动。\n\n"
                f"二、核心观点\n\n"
                f"基于宏观经济分析和行业研究，我们认为当前{topic_name.replace('报告', '').replace('分析', '')}领域"
                f"{random.choice(['存在结构性机会', '风险收益比较优', '需审慎对待', '配置价值凸显'])}。\n\n"
                f"三、投资建议\n\n"
                f"建议{random.choice(['积极配置', '适度参与', '防御为主', '精选个券'])}，"
                f"重点关注{random.choice(['优质龙头企业', '高股息标的', '景气度向上的细分赛道', '被低估的价值标的'])}。"
            ))
            attachments.append({"filename": fname, "type": "docx"})
        else:
            fname = f"research_{topic_key}_{i+1}.zip"
            att_path = str(ATTACHMENTS_DIR / fname)
            make_zip(att_path, {
                f"{topic_key}_report.md": (
                    f"# {topic_name}\n\n"
                    f"## Summary\n"
                    f"This research covers key developments in {topic_name}.\n\n"
                    f"## Data Sources\n"
                    f"- Wind Financial Terminal\n"
                    f"- CEIC\n"
                    f"- Company filings\n"
                    f"- Industry associations"
                ),
                f"{topic_key}_data.csv": (
                    "Date,Indicator,Value,YoY%,MoM%\n"
                    f"{date[:10]},GDP Growth,{macro_gdp},{round(macro_gdp - macro_gdp_prev, 1)},{round(random.uniform(-1,2), 1)}\n"
                    f"{date[:10]},CPI,{macro_cpi},{round(random.uniform(-1,3), 1)},{round(random.uniform(-0.5,1), 1)}\n"
                    f"{date[:10]},PMI,{macro_pmi},{round(random.uniform(-2,2), 1)},{round(random.uniform(-1,1), 1)}"
                ),
                f"disclaimer.txt": "This report is for internal use only. Not investment advice.",
            })
            attachments.append({"filename": fname, "type": "zip"})

        emails.append({
            "id": eid,
            "from": {"name": sender["name"], "email": sender["email"]},
            "to": [{"name": r["name"], "email": r["email"]} for r in recipients],
            "cc": [],
            "subject": subject,
            "body": body_text,
            "date": date,
            "tags": ["投研", topic_key],
            "attachments": attachments,
        })
    return emails


# ============================================================
# 主函数
# ============================================================

def main():
    random.seed(42)

    EMAILS_DIR.mkdir(parents=True, exist_ok=True)
    ATTACHMENTS_DIR.mkdir(parents=True, exist_ok=True)

    # 清理旧附件
    for f in ATTACHMENTS_DIR.iterdir():
        f.unlink()

    print("Generating financial email data...")

    all_emails = []
    offset = 1

    generators = [
        ("信托业务", gen_trust_emails, 15),
        ("基金/投资报告", gen_fund_emails, 15),
        ("合规/监管", gen_compliance_emails, 12),
        ("风控报告", gen_risk_emails, 10),
        ("法务/合同", gen_legal_emails, 12),
        ("运营管理", gen_operations_emails, 12),
        ("财富管理", gen_wealth_emails, 10),
        ("投研/市场分析", gen_research_emails, 14),
    ]

    for name, gen_func, count in generators:
        print(f"  Generating {name} ({count} emails)...")
        emails = gen_func(offset)
        all_emails.extend(emails)
        offset += count

    emails_file = EMAILS_DIR / "emails.json"
    with open(emails_file, "w", encoding="utf-8") as f:
        json.dump(all_emails, f, ensure_ascii=False, indent=2)

    print(f"\nDone! Generated {len(all_emails)} emails.")
    print(f"  Emails JSON: {emails_file}")
    print(f"  Attachments: {ATTACHMENTS_DIR}")

    att_types: dict[str, int] = {}
    for email in all_emails:
        for att in email["attachments"]:
            ext = att["type"]
            att_types[ext] = att_types.get(ext, 0) + 1
    print(f"\nAttachment stats:")
    for ext, count in sorted(att_types.items()):
        print(f"  .{ext}: {count} files")


if __name__ == "__main__":
    main()
