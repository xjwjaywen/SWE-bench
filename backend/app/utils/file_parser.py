"""Utility module for parsing attachment files and extracting text content."""

import csv
import io
import os
import zipfile
from pathlib import Path


def parse_pdf(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n".join(texts)
    except Exception:
        return ""


def parse_docx(file_path: str) -> str:
    try:
        from docx import Document

        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def parse_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            texts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    texts.append(line)
        wb.close()
        return "\n".join(texts)
    except Exception:
        return ""


def parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                texts.append(f"[Slide {i}]")
                texts.extend(slide_texts)
        return "\n".join(texts)
    except Exception:
        return ""


def parse_csv(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            lines = []
            for row in reader:
                line = "\t".join(row).strip()
                if line:
                    lines.append(line)
        return "\n".join(lines)
    except Exception:
        return ""


def parse_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:
        return ""


def parse_image(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    return f"[Image file: {os.path.basename(file_path)} ({ext.lstrip('.')} format)]"


def parse_zip(file_path: str) -> str:
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            names = zf.namelist()
            listing = "\n".join(f"  - {name}" for name in names)
            return f"[ZIP archive contents ({len(names)} files):\n{listing}\n]"
    except Exception:
        return ""


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    ".csv": parse_csv,
    ".txt": parse_txt,
    ".png": parse_image,
    ".jpg": parse_image,
    ".jpeg": parse_image,
    ".zip": parse_zip,
}


def parse_file(file_path: str) -> str:
    """Parse a file and extract text content based on its extension.

    Args:
        file_path: Path to the file to parse.

    Returns:
        Extracted text content, or empty string on failure.
    """
    ext = Path(file_path).suffix.lower()
    parser = PARSERS.get(ext)
    if parser is None:
        return ""
    try:
        return parser(file_path)
    except Exception:
        return ""
